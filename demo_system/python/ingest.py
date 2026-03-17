"""
Step 1 — Ingestion & Parsing
Walks all 4 data source folders, routes each file to the correct parser,
and outputs cache/parsed_docs.json.

Supported types:
  .eml   → email (headers + body)
  .txt   → plain text / meeting minutes / policy docs
  .xlsx  → Excel (per-sheet with headers)
  .csv   → CSV (with headers)
  .sas   → SAS script (plain text, preserves code)
  .log   → SAS log (plain text)
  .pptx  → PowerPoint (slide text)
"""

import csv
import email
import json
import re
import sys
from datetime import datetime, timezone
from email import policy as email_policy
from pathlib import Path

import openpyxl

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import pdfplumber
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    import docx as docx_lib
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

from config import (
    CACHE_DIR,
    DATA_SOURCES,
    EXCLUDE_DIRS,
    EXCLUDE_EXTENSIONS,
    EXCLUDE_FILES,
)
from metadata_rules import enrich_rule_based, reconstruct_email_threads


# ── Helpers ───────────────────────────────────────────────────────────────────

def now_iso() -> str:
    return datetime.now(timezone.utc).isoformat()


def make_doc_id(source_label: str, path: Path) -> str:
    """
    Emails:   stem is already 'MSG-001' or 'MTG-003' → use as-is
    Others:   {source_label}_{sanitised_stem}
    """
    stem = path.stem
    if source_label == "email":
        return stem  # MSG-001, MTG-003, MSG-012-V1 …
    safe = re.sub(r"[^a-zA-Z0-9]+", "_", stem).strip("_")
    return f"{source_label}_{safe}"


def classify_file_type(source_label: str, path: Path) -> str:
    ext = path.suffix.lower()
    if ext == ".eml":
        return "email"
    if ext == ".xlsx":
        return "excel_sheet"
    if ext == ".csv":
        return "csv"
    if ext in (".sas",):
        return "sas_script"
    if ext == ".log":
        return "log_file"
    if ext == ".pptx":
        return "powerpoint"
    if ext == ".pdf":
        return "pdf"
    if ext in (".docx", ".doc"):
        return "word_doc"
    if ext == ".txt":
        if source_label == "meetings":
            return "meeting_minutes"
        if source_label == "guide":
            return "policy_doc"
        return "word_doc"
    return "unknown"


def count_words(text: str) -> int:
    return len(text.split())


# ── Parsers ───────────────────────────────────────────────────────────────────

def parse_eml(path: Path) -> dict:
    with open(path, "rb") as f:
        msg = email.message_from_binary_file(f, policy=email_policy.default)

    def header(key):
        val = msg.get(key, "")
        return str(val).strip() if val else ""

    def header_list(key):
        val = header(key)
        if not val:
            return []
        return [a.strip() for a in val.split(",") if a.strip()]

    # Extract body
    body = ""
    attachments = []
    if msg.is_multipart():
        for part in msg.walk():
            ct = part.get_content_type()
            cd = str(part.get("Content-Disposition", ""))
            if "attachment" in cd:
                attachments.append(part.get_filename() or "unknown")
            elif ct == "text/plain":
                body += part.get_content() or ""
    else:
        body = msg.get_content() or ""

    structured = {
        "from":        header("From"),
        "to":          header_list("To"),
        "cc":          header_list("Cc"),
        "date":        header("Date"),
        "subject":     header("Subject"),
        "message_id":  header("Message-ID"),
        "in_reply_to": header("In-Reply-To") or None,
        "attachments": attachments,
    }

    # Build readable flat text for embedding
    to_str  = ", ".join(structured["to"])
    cc_str  = f"\nCc: {', '.join(structured['cc'])}" if structured["cc"] else ""
    att_str = f"\n[Attachments: {', '.join(attachments)}]" if attachments else ""

    raw_text = (
        f"From: {structured['from']}\n"
        f"To: {to_str}{cc_str}\n"
        f"Date: {structured['date']}\n"
        f"Subject: {structured['subject']}\n"
        f"{att_str}\n\n"
        f"{body.strip()}"
    )

    return {"raw_text": raw_text.strip(), "structured_fields": structured}


def parse_txt(path: Path) -> dict:
    text = path.read_text(encoding="utf-8", errors="replace").strip()
    return {"raw_text": text, "structured_fields": {}}


def parse_xlsx(path: Path) -> dict:
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    sheet_texts = []
    sheet_names = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue

        # First row = headers
        headers = [str(h) if h is not None else "" for h in rows[0]]
        lines = [" | ".join(headers)]

        for row in rows[1:]:
            cells = [str(c) if c is not None else "" for c in row]
            # Skip entirely empty rows
            if any(c.strip() for c in cells):
                lines.append(" | ".join(cells))

        sheet_text = f"Sheet: {sheet_name}\n" + "\n".join(lines)
        sheet_texts.append(sheet_text)
        sheet_names.append(sheet_name)

    wb.close()
    raw_text = "\n\n".join(sheet_texts)
    return {
        "raw_text": raw_text,
        "structured_fields": {"sheets": sheet_names},
    }


def parse_csv(path: Path) -> dict:
    lines = []
    with open(path, encoding="utf-8", errors="replace", newline="") as f:
        reader = csv.reader(f)
        rows = list(reader)

    if not rows:
        return {"raw_text": "", "structured_fields": {}}

    headers = rows[0]
    lines.append(" | ".join(headers))
    for row in rows[1:]:
        if any(cell.strip() for cell in row):
            lines.append(" | ".join(row))

    raw_text = f"File: {path.name}\n" + "\n".join(lines)
    return {
        "raw_text": raw_text,
        "structured_fields": {"columns": headers, "row_count": len(rows) - 1},
    }


def parse_sas(path: Path) -> dict:
    text = path.read_text(encoding="utf-8", errors="replace").strip()
    return {"raw_text": text, "structured_fields": {"language": "SAS"}}


def parse_log(path: Path) -> dict:
    text = path.read_text(encoding="utf-8", errors="replace").strip()
    return {"raw_text": text, "structured_fields": {"log_type": "SAS_execution"}}


def parse_pptx(path: Path) -> dict:
    if not PPTX_AVAILABLE:
        return {
            "raw_text": f"[PPTX parser unavailable — install python-pptx to parse {path.name}]",
            "structured_fields": {},
        }
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
        if texts:
            slides.append(f"--- Slide {i} ---\n" + "\n".join(texts))

    raw_text = "\n\n".join(slides)
    return {
        "raw_text": raw_text,
        "structured_fields": {"slide_count": len(prs.slides)},
    }


def parse_pdf(path: Path) -> dict:
    if not PDF_AVAILABLE:
        raise ImportError("pdfplumber not installed — run: pip install pdfplumber")
    pages = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            if text.strip():
                pages.append(text.strip())
    return {
        "raw_text": "\n\n".join(pages),
        "structured_fields": {"page_count": len(pages)},
    }


def parse_docx(path: Path) -> dict:
    if not DOCX_AVAILABLE:
        raise ImportError("python-docx not installed — run: pip install python-docx")
    doc = docx_lib.Document(str(path))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return {
        "raw_text": "\n\n".join(paragraphs),
        "structured_fields": {"paragraph_count": len(paragraphs)},
    }


def parse_word_doc(path: Path) -> dict:
    """Route .docx/.doc to python-docx, .txt to plain text parser."""
    if path.suffix.lower() in (".docx", ".doc"):
        return parse_docx(path)
    return parse_txt(path)


# ── Router ────────────────────────────────────────────────────────────────────

PARSERS = {
    "email":           parse_eml,
    "meeting_minutes": parse_txt,
    "word_doc":        parse_word_doc,
    "policy_doc":      parse_txt,
    "excel_sheet":     parse_xlsx,
    "csv":             parse_csv,
    "sas_script":      parse_sas,
    "log_file":        parse_log,
    "powerpoint":      parse_pptx,
    "pdf":             parse_pdf,
}


def parse_file(source_label: str, path: Path) -> dict:
    doc_id    = make_doc_id(source_label, path)
    file_type = classify_file_type(source_label, path)
    timestamp = now_iso()

    try:
        parser = PARSERS.get(file_type)
        if parser is None:
            raise ValueError(f"No parser for file_type '{file_type}'")
        result = parser(path)
        raw_text = result["raw_text"]
        status   = "success" if raw_text.strip() else "partial"
        error    = None if status == "success" else "Empty content after parsing"
    except Exception as exc:
        raw_text = ""
        result   = {"raw_text": "", "structured_fields": {}}
        status   = "failed"
        error    = str(exc)

    return {
        "doc_id":           doc_id,
        "source_path":      str(path),
        "source_folder":    source_label,
        "file_type":        file_type,
        "file_extension":   path.suffix.lower(),
        "raw_text":         raw_text,
        "structured_fields": result.get("structured_fields", {}),
        "parse_status":     status,
        "parse_error":      error,
        "parse_timestamp":  timestamp,
        "char_count":       len(raw_text),
        "word_count":       count_words(raw_text),
    }


# ── Walker ────────────────────────────────────────────────────────────────────

def collect_files(source_label: str, folder: Path) -> list[Path]:
    files = []
    for path in sorted(folder.rglob("*")):
        if path.is_dir():
            continue
        if any(ex in path.parts for ex in EXCLUDE_DIRS):
            continue
        if path.suffix.lower() in EXCLUDE_EXTENSIONS:
            continue
        if path.name in EXCLUDE_FILES:
            continue
        files.append(path)
    return files


# ── Main ──────────────────────────────────────────────────────────────────────

def run_ingestion(verbose: bool = True) -> list[dict]:
    CACHE_DIR.mkdir(parents=True, exist_ok=True)
    all_docs = []

    for source_label, folder in DATA_SOURCES.items():
        if not folder.exists():
            if verbose:
                print(f"  [SKIP] {source_label}: folder not found at {folder}")
            continue

        files = collect_files(source_label, folder)
        if verbose:
            print(f"\n[{source_label.upper()}] {len(files)} files found in {folder}")

        for path in files:
            doc = parse_file(source_label, path)
            doc = enrich_rule_based(doc)
            all_docs.append(doc)
            if verbose:
                icon = "✓" if doc["parse_status"] == "success" else \
                       "~" if doc["parse_status"] == "partial" else "✗"
                print(f"  {icon} {doc['doc_id']:40s}  "
                      f"{doc['file_type']:20s}  "
                      f"{doc['word_count']:>5} words  "
                      f"[{doc['parse_status']}]")
                if doc["parse_error"]:
                    print(f"      ERROR: {doc['parse_error']}")

    # Second pass: wire up email thread IDs across all docs
    all_docs = reconstruct_email_threads(all_docs)

    out_path = CACHE_DIR / "parsed_docs.json"
    with open(out_path, "w", encoding="utf-8") as f:
        json.dump(all_docs, f, indent=2, ensure_ascii=False)

    if verbose:
        success = sum(1 for d in all_docs if d["parse_status"] == "success")
        partial = sum(1 for d in all_docs if d["parse_status"] == "partial")
        failed  = sum(1 for d in all_docs if d["parse_status"] == "failed")
        print(f"\n{'─'*60}")
        print(f"Total parsed : {len(all_docs)}")
        print(f"  ✓ success  : {success}")
        print(f"  ~ partial  : {partial}")
        print(f"  ✗ failed   : {failed}")
        print(f"Output       : {out_path}")

    return all_docs


def load_parsed_docs() -> list[dict]:
    path = CACHE_DIR / "parsed_docs.json"
    if not path.exists():
        return []
    with open(path, encoding="utf-8") as f:
        return json.load(f)


def get_ingestion_status() -> dict:
    docs = load_parsed_docs()
    if not docs:
        return {"status": "not_run", "total": 0}

    by_type    = {}
    by_folder  = {}
    failed     = []

    for d in docs:
        by_type[d["file_type"]]      = by_type.get(d["file_type"], 0) + 1
        by_folder[d["source_folder"]] = by_folder.get(d["source_folder"], 0) + 1
        if d["parse_status"] == "failed":
            failed.append({"doc_id": d["doc_id"], "error": d["parse_error"]})

    return {
        "status":       "complete",
        "total":        len(docs),
        "success":      sum(1 for d in docs if d["parse_status"] == "success"),
        "partial":      sum(1 for d in docs if d["parse_status"] == "partial"),
        "failed":       len(failed),
        "failed_docs":  failed,
        "by_file_type": by_type,
        "by_folder":    by_folder,
        "last_run":     docs[0]["parse_timestamp"] if docs else None,
    }


if __name__ == "__main__":
    verbose = "--quiet" not in sys.argv
    run_ingestion(verbose=verbose)
